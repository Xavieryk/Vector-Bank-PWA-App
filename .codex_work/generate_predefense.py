from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DPt
from PIL import Image, ImageDraw

OUT_DIR = Path("C:/Users/Admin/Desktop/д")
SCREEN_DIR = OUT_DIR / "22"
WORK_DIR = Path("X:/PWAAPP/.codex_work")
PPTX_OUT = OUT_DIR / "predefense_mobile_bank_pwa.pptx"
DOCX_OUT = OUT_DIR / "doklad_mobile_bank_pwa.docx"
MD_OUT = OUT_DIR / "doklad_mobile_bank_pwa.md"
BG_DIR = WORK_DIR / "deck_assets"
BG_DIR.mkdir(parents=True, exist_ok=True)

NAVY = "071724"
PANEL = "1E2B34"
ACCENT = "ABFF1A"
TEXT = "DEE3EA"
DANGER = "FF2024"
MUTED = "95A4B2"
FONT = "Poppins"
SLIDE_W, SLIDE_H = 13.333, 7.5


def rgb(hexstr: str) -> RGBColor:
    h = hexstr.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def make_bg() -> Path:
    w, h = 1920, 1080
    img = Image.new("RGB", (w, h), "#071724")
    px = img.load()
    for y in range(h):
        for x in range(w):
            t = y / h
            v = int(13 * t)
            px[x, y] = (7 + v // 4, 23 + v // 2, 36 + v)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    for x in range(-w, w * 2, 64):
        d.line((x, 0, x + 420, h), fill=(30, 43, 52, 42), width=2)
    d.ellipse((-260, 620, 660, 1480), fill=(8, 24, 37, 95))
    d.ellipse((1280, -380, 2180, 520), fill=(171, 255, 26, 22))
    img = Image.alpha_composite(img.convert("RGBA"), overlay)
    path = BG_DIR / "deck-bg.png"
    img.save(path)
    return path


BG_PATH = make_bg()
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
blank = prs.slide_layouts[6]


def add_bg(slide):
    slide.shapes.add_picture(str(BG_PATH), 0, 0, width=prs.slide_width, height=prs.slide_height)


def set_text(tf, text, size=24, color=TEXT, bold=False, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.bold = bold


def add_text(slide, x, y, w, h, text, size=24, color=TEXT, bold=False, align=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_text(tf, text, size=size, color=color, bold=bold, align=align)
    if line_spacing:
        for p in tf.paragraphs:
            p.line_spacing = line_spacing
    return box


def add_panel(slide, x, y, w, h, fill=PANEL, radius=True, line=None, alpha=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if alpha is not None:
        shp.fill.transparency = alpha
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_chip(slide, x, y, text, fill=PANEL, color=TEXT, size=13, w=None, accent=False):
    width = w if w else max(1.0, 0.18 * len(text) + 0.35)
    add_panel(slide, x, y, width, 0.36, fill=fill, radius=True, line=ACCENT if accent else None)
    add_text(slide, x + 0.13, y + 0.075, width - 0.26, 0.2, text, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_title(slide, title, subtitle=None, num=None):
    if num:
        add_chip(slide, 0.62, 0.48, f"{num:02}", fill=ACCENT, color=NAVY, size=12, w=0.58)
        tx = 1.35
    else:
        tx = 0.72
    add_text(slide, tx, 0.42, 8.6, 0.7, title, size=30, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, tx, 1.05, 8.8, 0.36, subtitle, size=13.5, color=MUTED)


def add_bullets(slide, x, y, w, bullets, size=16, color=TEXT, accent=ACCENT, gap=0.62):
    cur = y
    for b in bullets:
        add_panel(slide, x, cur + 0.055, 0.11, 0.11, fill=accent, radius=True)
        add_text(slide, x + 0.28, cur, w - 0.28, 0.48, b, size=size, color=color)
        cur += gap


def add_metric(slide, x, y, value, label, accent_color=ACCENT):
    add_panel(slide, x, y, 2.25, 1.22, fill=PANEL, radius=True)
    add_text(slide, x + 0.2, y + 0.18, 1.85, 0.36, value, size=25, color=accent_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, y + 0.66, 1.85, 0.36, label, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_phone_image(slide, img_path: Path, x, y, h, label=None):
    img = Image.open(img_path)
    iw, ih = img.size
    width = h * iw / ih
    add_panel(slide, x - 0.06, y - 0.06, width + 0.12, h + 0.12, fill="0A1C2A", radius=True, line="2B3A45")
    slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), height=Inches(h))
    if label:
        add_chip(slide, x + 0.14, y + h - 0.48, label, fill="081825", color=TEXT, size=9.5, w=max(1.1, width - 0.28))
    return width


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2.2)


def screen(name: str) -> Path:
    return SCREEN_DIR / name


def new_slide(title=None, subtitle=None, num=None):
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    if title:
        add_title(slide, title, subtitle, num)
    return slide


slides_data = []

# 1
s = new_slide()
add_panel(s, 0.68, 0.58, 12.0, 6.35, fill=NAVY, radius=True, line="21323D")
add_chip(s, 0.95, 0.92, "ВКР · предзащита 2026", fill=PANEL, color=TEXT, w=2.35)
add_text(s, 0.95, 1.55, 7.1, 2.2, "Проектирование главного экрана мобильного банковского приложения", size=34, color=TEXT, bold=True, line_spacing=0.9)
add_text(s, 0.98, 4.25, 3.9, 0.28, "Студент: Зуев Никита Игоревич", size=14, color=TEXT, bold=True)
add_text(s, 0.98, 4.72, 3.9, 0.28, "Группа: 8.209-2", size=13, color=MUTED)
add_text(s, 0.98, 5.12, 5.0, 0.28, "Научный руководитель: Коняева Софья Игоревна", size=13, color=MUTED)
add_text(s, 0.98, 6.18, 3.0, 0.28, "PWA-прототип на React", size=13, color=ACCENT, bold=True)
add_phone_image(s, screen("главный экран.png"), 9.02, 0.72, 5.95, label="главный экран")
slides_data.append(("Титульный слайд", "Добрый день. Меня зовут Зуев Никита Игоревич. Тема моей выпускной квалификационной работы — проектирование главного экрана мобильного банковского приложения. В рамках работы я не ограничился макетами: результатом стал рабочий PWA-прототип на React, который можно открыть на мобильном устройстве и пройти основные банковские сценарии."))

# 2
s = new_slide("Актуальность", "Почему главный экран банка становится критичным интерфейсом", 2)
add_metric(s, 0.85, 1.68, "PWA", "альтернатива установке приложения")
add_metric(s, 3.35, 1.68, "Daily", "быстрые финансовые сценарии")
add_metric(s, 5.85, 1.68, "0 ошибок", "цель для денежных операций")
add_panel(s, 8.72, 1.25, 3.72, 4.92, fill=PANEL, radius=True)
add_text(s, 9.05, 1.6, 3.05, 0.48, "Контекст", size=20, color=TEXT, bold=True)
add_bullets(s, 9.05, 2.25, 3.05, [
    "банк открывают для коротких действий",
    "ошибка в сумме или получателе имеет цену",
    "санкции и ограничения установки усиливают роль web/PWA",
    "главный экран должен давать контроль, а не шум",
], size=13.3, gap=0.72)
add_text(s, 0.85, 3.45, 6.95, 1.3, "В дипломе главный экран рассматривается не как декоративная витрина, а как точка входа в ежедневные финансовые сценарии: баланс, история, быстрые действия и переводы.", size=21, color=TEXT, bold=True)
add_text(s, 0.9, 5.35, 6.5, 0.55, "Основной вызов: совместить скорость, безопасность и персонализацию без перегрузки интерфейса.", size=15, color=MUTED)
slides_data.append(("Актуальность", "Актуальность связана с тем, что мобильный банк стал ежедневным инструментом. Человек открывает его, чтобы быстро проверить баланс, увидеть последние операции или отправить перевод. При этом интерфейс работает с деньгами, поэтому любая ошибка восприятия суммы, получателя или статуса операции может привести к реальным последствиям. Дополнительно важен контекст ограничений установки банковских приложений, особенно на iPhone, из-за чего PWA и webview становятся практичным направлением."))

# 3
s = new_slide("Цель и рамка работы", "Что именно проектировалось и проверялось", 3)
add_panel(s, 0.82, 1.45, 5.9, 1.35, fill=PANEL, radius=True, line=ACCENT)
add_text(s, 1.12, 1.72, 5.25, 0.76, "Разработать и реализовать минимально работающий PWA-прототип главного экрана мобильного банка с персонализируемыми виджетами и базовыми сценариями.", size=15.5, color=TEXT, bold=True)
add_panel(s, 7.15, 1.45, 5.15, 1.35, fill=PANEL, radius=True)
add_text(s, 7.45, 1.72, 4.55, 0.76, "Объект: мобильное банковское приложение как цифровой сервис ежедневных финансовых операций.", size=15, color=TEXT)
add_panel(s, 7.15, 3.1, 5.15, 1.25, fill=PANEL, radius=True)
add_text(s, 7.45, 3.36, 4.55, 0.62, "Предмет: интерфейсные и технические решения главного экрана, баланса, переводов, истории и виджетов.", size=15, color=TEXT)
add_text(s, 0.88, 3.38, 2.4, 0.4, "Ключевые задачи", size=19, color=TEXT, bold=True)
add_bullets(s, 0.9, 4.03, 5.5, [
    "проанализировать аналоги и UX-критерии",
    "выявить пользовательские боли и daily-сценарии",
    "спроектировать навигацию, виджеты и сценарий перевода",
    "реализовать React/PWA и LocalStorage",
    "проверить работоспособность и юзабилити",
], size=13.5, gap=0.56)
add_text(s, 7.22, 5.18, 4.95, 0.82, "Фокус защиты: показать не только красивый интерфейс, но и цепочку решений от исследования до работающего прототипа.", size=17, color=ACCENT, bold=True)
slides_data.append(("Цель и рамка работы", "Цель работы — разработать и реализовать минимально работающий PWA-прототип главного экрана мобильного банка. Объектом является мобильное банковское приложение как сервис ежедневных финансовых операций. Предмет — интерфейсные и технические решения главного экрана: баланс, виджеты, переводы, история операций и персонализация. Задачи выстроены от анализа и требований до реализации и юзабилити-проверки."))

# 4
s = new_slide("Методы", "Исследовательская и проектная логика", 4)
steps = [
    ("1", "Анализ источников", "UX Problems Guide, стандарты, литература"),
    ("2", "Бенчмаркинг", "российские и зарубежные банки"),
    ("3", "Сценарный анализ", "daily-задачи и corner cases"),
    ("4", "Прототипирование", "Figma → React PWA"),
    ("5", "Юзабилити-проверка", "cognitive walkthrough, self-test"),
]
x = 0.8
for i, (n, t, dsc) in enumerate(steps):
    add_panel(s, x, 2.08, 2.15, 2.15, fill=PANEL, radius=True)
    add_chip(s, x + 0.18, 2.28, n, fill=ACCENT, color=NAVY, size=15, w=0.48)
    add_text(s, x + 0.18, 2.92, 1.75, 0.5, t, size=16.5, color=TEXT, bold=True)
    add_text(s, x + 0.18, 3.55, 1.75, 0.42, dsc, size=10.7, color=MUTED)
    if i < len(steps) - 1:
        add_arrow(s, x + 2.18, 3.15, x + 2.48, 3.15)
    x += 2.48
add_panel(s, 1.05, 5.12, 11.05, 0.92, fill="0A1C2A", radius=True, line="263844")
add_text(s, 1.35, 5.38, 10.4, 0.32, "Методы нужны не “для списка”, а чтобы обосновать интерфейсные решения: где нужен акцент, где предупреждение, где сохранение состояния, а где fallback.", size=15.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
slides_data.append(("Методы", "В работе использованы анализ литературы и открытых UX-исследований, сравнительный анализ аналогов, бенчмаркинг, сценарный анализ, прототипирование, cognitive walkthrough и техническое тестирование PWA. Эта логика важна: сначала я изучаю контекст и проблемы, затем превращаю их в требования, после этого собираю интерфейс и проверяю его на сценариях."))

# 5
s = new_slide("Пользовательские проблемы", "Что должен решить главный экран", 5)
add_panel(s, 0.8, 1.55, 5.55, 4.95, fill=PANEL, radius=True)
add_text(s, 1.12, 1.9, 4.8, 0.4, "Основные боли", size=21, color=TEXT, bold=True)
insights = [
    ("Страх ошибки", "сумма, получатель, карта списания"),
    ("Перегрузка", "слишком много функций на одном экране"),
    ("Невидимые состояния", "непонятно: отправлено, сохранено, требуется ли подтверждение"),
    ("Низкий контроль", "сложно быстро проверить баланс и историю"),
]
y = 2.55
for title, desc in insights:
    add_panel(s, 1.15, y + 0.06, 0.16, 0.16, fill=ACCENT, radius=True)
    add_text(s, 1.5, y, 2.7, 0.28, title, size=15.5, color=TEXT, bold=True)
    add_text(s, 3.42, y, 2.25, 0.28, desc, size=11.8, color=MUTED)
    y += 0.78
add_panel(s, 6.85, 1.55, 5.6, 4.95, fill="0A1C2A", radius=True, line="243744")
add_text(s, 7.18, 1.9, 4.9, 0.4, "Проектные ответы", size=21, color=ACCENT, bold=True)
add_bullets(s, 7.18, 2.6, 4.8, [
    "подтверждение крупной суммы",
    "история и детали операции",
    "ошибка “недостаточно средств” как модальное окно",
    "виджеты с настройкой и сохранением",
    "явные toast-сообщения и состояния",
], size=14, gap=0.62)
slides_data.append(("Пользовательские проблемы", "По результатам анализа и экспертных проверок ключевые проблемы связаны не только с визуальной эстетикой. В банковском интерфейсе важны контроль и понятность состояний. Пользователь боится ошибиться с получателем или суммой, может не понимать, сохранилось ли действие, и быстро теряется, если экран перегружен. Поэтому проектные решения направлены на предотвращение ошибок: подтверждение крупного перевода, явные состояния, детали операций, сохранение настроек и обработка недостатка средств."))

# 6
s = new_slide("Глава 1", "Аналитическая основа проекта", 6)
add_text(s, 0.86, 1.55, 5.2, 0.62, "Анализ предметной области и требований к главному экрану мобильного банка", size=25, color=TEXT, bold=True)
add_bullets(s, 0.9, 2.62, 5.6, [
    "роль главного экрана: баланс, действия, история",
    "UX-критерии: понятность, обратная связь, доступность, предотвращение ошибок",
    "аналоги: Сбер, Т-Банк, Альфа-Банк, ВТБ, Revolut, Monzo, N26, Wise",
    "требования: виджеты, LocalStorage, подтверждение крупных переводов",
], size=14.2, gap=0.68)
add_panel(s, 7.05, 1.44, 4.95, 4.95, fill=PANEL, radius=True)
add_text(s, 7.38, 1.85, 4.25, 0.4, "Вывод главы", size=22, color=ACCENT, bold=True)
add_text(s, 7.4, 2.52, 4.15, 1.25, "Главный экран должен быть персонализируемым, но не терять банковскую ясность: пользователь сначала видит финансовое состояние, затем быстрые действия и подтверждение важных операций.", size=17, color=TEXT, bold=True)
add_text(s, 7.4, 4.35, 4.05, 0.7, "Критерий качества: сценарий должен завершаться без догадок и без риска случайного действия.", size=15, color=MUTED)
slides_data.append(("Глава 1", "В первой главе рассматривается предметная область: роль главного экрана, критерии UX-качества, аналогичные решения российских и зарубежных банков, пользовательские боли и требования. Главный вывод: экран должен не просто красиво показывать карту, а помогать быстро понять финансовое состояние и безопасно перейти к действиям."))

# 7
s = new_slide("Глава 2", "Проектирование экранов и сценариев", 7)
for name, x, label in [
    ("переводы.png", 0.82, "переводы"),
    ("перевод по ном. тел. СБП.png", 3.35, "сумма"),
    ("при большом платеже.png", 5.88, "проверка"),
    ("потверждение-1.png", 8.41, "код"),
    ("потверждение-3.png", 10.94, "результат"),
]:
    add_phone_image(s, screen(name), x, 1.25, 4.7, label=label)
add_text(s, 0.88, 6.35, 11.8, 0.35, "Сценарий перевода: выбор получателя → ввод суммы → проверка риска → подтверждение → результат и чек", size=15.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
slides_data.append(("Глава 2", "Вторая глава посвящена проектированию: концепции продукта, информационной архитектуре, навигации, главному экрану, виджетам, сценарию перевода и UI-kit. На слайде показан один из ключевых пользовательских путей: пользователь выбирает получателя, вводит сумму, проверяет детали, подтверждает крупный перевод и получает результат."))

# 8
s = new_slide("Проектный результат", "Рабочий прототип вместо статичного макета", 8)
add_phone_image(s, screen("главный экран.png"), 0.9, 1.25, 5.35, label="главный экран")
add_panel(s, 4.15, 1.46, 7.95, 4.98, fill=PANEL, radius=True)
add_text(s, 4.55, 1.88, 7.15, 0.45, "Что реализовано", size=23, color=TEXT, bold=True)
add_bullets(s, 4.58, 2.62, 6.9, [
    "главный экран с картами и общим балансом",
    "2–3 настраиваемых виджета с LocalStorage",
    "перевод по номеру телефона и выбор банка",
    "история операций, детали и повтор перевода",
    "шаблоны, избранное и экран “Все”",
    "ошибки, чек, QR-сканер, профиль и настройки",
], size=13.8, gap=0.52)
add_chip(s, 4.6, 5.7, "React", fill="0A1C2A", color=TEXT, w=1.0)
add_chip(s, 5.8, 5.7, "Vite", fill="0A1C2A", color=TEXT, w=0.9)
add_chip(s, 6.9, 5.7, "PWA", fill=ACCENT, color=NAVY, w=0.9)
add_chip(s, 8.0, 5.7, "LocalStorage", fill="0A1C2A", color=TEXT, w=1.55)
add_chip(s, 9.75, 5.7, "Web API", fill="0A1C2A", color=TEXT, w=1.2)
slides_data.append(("Проектный результат", "Итоговый проект — это рабочий PWA-прототип. В нём есть главный экран с картами и общим балансом, настраиваемые виджеты, перевод по номеру телефона, выбор карты списания и банка получателя, история операций, детали, повтор перевода, шаблоны, QR-сканер, профиль и настройки. Это соответствует требованию диплома: результат не только макетный, а интерактивный и демонстрируемый на телефоне."))

# 9
s = new_slide("Главный экран", "Персонализация без потери банковской ясности", 9)
for x, title, bullets in [
    (0.85, "Карточки", ["свайп вверх/вниз", "скрытие баланса по карте", "общий баланс как агрегатор"]),
    (4.62, "Виджеты", ["последние операции", "быстрые действия", "курсы валют", "скрытие и перестановка"]),
    (8.39, "Настройки", ["сохранение на устройстве", "порог подтверждения", "toast обратной связи"]),
]:
    add_panel(s, x, 1.38, 3.42, 4.95, fill=PANEL, radius=True)
    add_text(s, x + 0.3, 1.72, 2.8, 0.38, title, size=20, color=TEXT, bold=True)
    add_bullets(s, x + 0.3, 2.35, 2.85, bullets, size=13, gap=0.62)
add_text(s, 0.96, 6.55, 10.95, 0.26, "Требование ВКР “2–3 виджета, перестановка/скрытие и LocalStorage” закрыто в прототипе.", size=14.5, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
slides_data.append(("Главный экран", "Главный экран решает две задачи: быстро показать финансовое состояние и дать доступ к действиям. Карты можно перелистывать, баланс можно скрывать по конкретной карте, общий баланс агрегирует остальные счета. Виджеты можно скрывать и переставлять, а настройки сохраняются на устройстве через LocalStorage."))

# 10
s = new_slide("Сценарий перевода", "Безопасность и обратная связь", 10)
for name, x, label in [
    ("перевод по ном. тел. СБП.png", 0.92, "ввод суммы"),
    ("при большом платеже.png", 3.68, "проверка"),
    ("потверждение-4.png", 6.44, "ошибка кода"),
    ("потверждение-3.png", 9.2, "успех"),
]:
    add_phone_image(s, screen(name), x, 1.25, 5.15, label=label)
add_panel(s, 0.95, 6.46, 10.9, 0.52, fill="081825", radius=True, line="263844")
add_text(s, 1.25, 6.62, 10.3, 0.18, "Крупная сумма проходит через усиленное подтверждение; недостаток средств и неверный код обрабатываются отдельными состояниями.", size=12.7, color=TEXT, align=PP_ALIGN.CENTER)
slides_data.append(("Сценарий перевода", "Сценарий перевода спроектирован вокруг снижения риска. Пользователь видит получателя, карту списания, банк, сумму и комиссию. Если сумма достигает заданного порога, появляется усиленное подтверждение. Если средств недостаточно, это показывается во всплывающем окне, не ломая путь пользователя. После успешного перевода деньги списываются с карты, операция попадает в историю, а пользователь может открыть чек или повторить перевод."))

# 11
s = new_slide("ИТ-компонент", "Что демонстрирует прототип технически", 11)
add_panel(s, 0.85, 1.42, 5.25, 4.75, fill=PANEL, radius=True)
add_text(s, 1.18, 1.82, 4.55, 0.4, "Стек и функции", size=22, color=TEXT, bold=True)
add_bullets(s, 1.18, 2.45, 4.55, [
    "React-компоненты и состояния",
    "адаптивная верстка под мобильные экраны",
    "PWA manifest и запуск как web-приложение",
    "LocalStorage для настроек и операций",
    "Web API: контакты/камера с fallback-сценариями",
], size=13.6, gap=0.55)
add_panel(s, 6.65, 1.42, 5.85, 4.75, fill="0A1C2A", radius=True, line="263844")
add_text(s, 7.0, 1.82, 4.9, 0.4, "Компонентная логика", size=22, color=ACCENT, bold=True)
for label, x, y in [("UI", 7.05, 2.62), ("State", 8.7, 2.62), ("Storage", 10.35, 2.62), ("PWA", 8.7, 4.2)]:
    add_panel(s, x, y, 1.28, 0.78, fill=PANEL, radius=True, line=ACCENT if label == "PWA" else None)
    add_text(s, x + 0.08, y + 0.24, 1.1, 0.2, label, size=14.5, color=TEXT if label != "PWA" else ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_arrow(s, 8.35, 3.0, 8.68, 3.0)
add_arrow(s, 9.98, 3.0, 10.35, 3.0)
add_arrow(s, 9.32, 3.43, 9.32, 4.18)
add_text(s, 7.05, 5.25, 4.95, 0.35, "Код можно показать как доказательство ИТ-составляющей: JSX/TS, CSS, manifest, логика состояний.", size=13, color=MUTED, align=PP_ALIGN.CENTER)
slides_data.append(("ИТ-компонент", "ИТ-компонент работы — это React/PWA-прототип. В нём используются компонентная структура, состояния интерфейса, адаптивная верстка, LocalStorage для настроек и истории, PWA-конфигурация, а также Web API: контакты и камера для QR-сканера с fallback-сообщениями, если браузер или устройство ограничивают доступ."))

# 12
s = new_slide("Юзабилити-тестирование", "Что проверялось и что было исправлено", 12)
for x, title, color, bullets, accent in [
    (0.82, "Метод", TEXT, ["cognitive walkthrough", "self usability testing", "4 ролевые модели", "12 ключевых сценариев"], ACCENT),
    (4.72, "Найдено", DANGER, ["перегрузка переводов", "неочевидность настроек", "toast и состояния", "длинные имена", "ограничения QR API"], DANGER),
    (8.62, "Исправлено", ACCENT, ["toast-обратная связь", "fallback-сценарии", "ограничение длинных имён", "экран “Все” для шаблонов", "модальное “недостаточно средств”"], ACCENT),
]:
    add_panel(s, x, 1.35, 3.65, 4.95, fill=PANEL, radius=True)
    add_text(s, x + 0.3, 1.75, 3.05, 0.4, title, size=21, color=color, bold=True)
    add_bullets(s, x + 0.3, 2.4, 3.05, bullets, size=13.5, gap=0.62, accent=accent)
add_text(s, 1.0, 6.58, 11.0, 0.23, "Предварительная проверка: критических блокирующих проблем не выявлено; внешнее тестирование N=5 запланировано до 10.05.2026.", size=12.8, color=MUTED, align=PP_ALIGN.CENTER)
slides_data.append(("Юзабилити-тестирование", "Предварительно проведены cognitive walkthrough и self usability testing. Проверялись 12 сценариев: главный экран, перевод, крупная сумма, недостаток средств, чек, история, повтор перевода, шаблоны, настройки, QR и профиль. Были выявлены риски: перегрузка экрана переводов, неочевидный вход в настройки, необходимость toast-сообщений, длинные имена и зависимость QR от браузера. Часть проблем уже исправлена, а внешнее тестирование на пяти респондентах запланировано до 10 мая 2026 года."))

# 13
s = new_slide("Текущая готовность", "Что уже есть и что осталось до защиты", 13)
add_panel(s, 0.88, 1.38, 5.45, 5.05, fill=PANEL, radius=True, line=ACCENT)
add_text(s, 1.22, 1.8, 4.8, 0.4, "Готово", size=22, color=ACCENT, bold=True)
add_bullets(s, 1.22, 2.45, 4.7, [
    "черновик ВКР 50+ страниц",
    "56 источников",
    "ключевые макеты и UI-kit",
    "React/PWA-прототип",
    "LocalStorage, QR, контакты, история",
    "предварительный self-test",
], size=13.6, gap=0.52)
add_panel(s, 6.95, 1.38, 5.45, 5.05, fill=PANEL, radius=True)
add_text(s, 7.29, 1.8, 4.8, 0.4, "Осталось", size=22, color=TEXT, bold=True)
add_bullets(s, 7.29, 2.45, 4.7, [
    "внешнее тестирование N=5",
    "обновить выводы по результатам",
    "сверить библиографию по ГОСТ",
    "подготовить финальные скриншоты",
    "проверить демонстрационный телефон и QR",
], size=13.6, gap=0.52, accent=MUTED)
add_text(s, 1.0, 6.62, 11.0, 0.23, "Финальная версия работы и документы к защите: до 18 мая 2026 г.", size=13.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
slides_data.append(("Текущая готовность", "На текущий момент готов черновик работы более 50 страниц, список из 56 источников, структура исследования, макеты, UI-kit, рабочий PWA-прототип и предварительное юзабилити-тестирование. Осталось провести внешнее тестирование на пяти респондентах, обновить выводы, сверить источники по ГОСТ, подготовить финальные скриншоты и заранее проверить демонстрационный телефон и QR-сценарий."))

# 14
s = new_slide("Итог", "Главный экран как система контроля, а не просто витрина", 14)
add_text(s, 0.9, 1.45, 6.55, 1.35, "Результат работы — связка исследования, дизайна и реализации: от UX-проблем к интерактивному банковскому PWA.", size=28, color=TEXT, bold=True)
for x, title, sub, fill, color in [
    (0.95, "Исследование", "проблемы и критерии", PANEL, TEXT),
    (4.95, "Дизайн", "сценарии и UI-kit", PANEL, TEXT),
    (8.95, "PWA", "рабочий прототип", ACCENT, NAVY),
]:
    add_panel(s, x, 3.35, 3.1, 1.55, fill=fill, radius=True)
    add_text(s, x + 0.25, 3.72, 2.55, 0.3, title, size=19, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.25, 4.22, 2.55, 0.28, sub, size=12, color=color if fill == ACCENT else MUTED, align=PP_ALIGN.CENTER)
add_text(s, 0.98, 6.12, 10.85, 0.35, "Спасибо за внимание. Готов показать прототип и ответить на вопросы.", size=19, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
slides_data.append(("Итог", "В итоге работа показывает полный проектный цикл: исследование, проектирование, реализацию и проверку. Главный экран в моей работе — это не просто красивая карточка банка, а система контроля: пользователь видит баланс, понимает последние операции, может настроить виджеты и безопасно перейти к переводу. Спасибо за внимание. Я готов показать прототип и ответить на вопросы."))

for idx, slide in enumerate(prs.slides, start=1):
    if idx != 1:
        add_text(slide, 12.05, 6.95, 0.55, 0.18, str(idx), size=8.5, color="52616B", align=PP_ALIGN.RIGHT)

prs.save(PPTX_OUT)

doc = Document()
doc.styles["Normal"].font.name = FONT
doc.styles["Normal"].font.size = DPt(11)
doc.styles["Title"].font.name = FONT
doc.styles["Title"].font.size = DPt(18)
doc.styles["Heading 1"].font.name = FONT
doc.styles["Heading 1"].font.size = DPt(14)

doc.add_heading("Доклад к презентации: проектирование главного экрана мобильного банковского приложения", 0)
p = doc.add_paragraph()
p.add_run("Студент: ").bold = True
p.add_run("Зуев Никита Игоревич, группа 8.209-2")
p = doc.add_paragraph()
p.add_run("Научный руководитель: ").bold = True
p.add_run("Коняева Софья Игоревна")
p = doc.add_paragraph()
p.add_run("Рекомендуемый темп: ").bold = True
p.add_run("6–8 минут. Слайды служат опорой, основной текст произносится устно.")

for i, (title, speech) in enumerate(slides_data, start=1):
    doc.add_heading(f"Слайд {i}. {title}", level=1)
    doc.add_paragraph(speech)

doc.add_heading("Короткая версия на случай ограничения по времени", level=1)
doc.add_paragraph(
    "Добрый день. Тема моей ВКР — проектирование главного экрана мобильного банковского приложения. "
    "Актуальность связана с тем, что мобильный банк стал ежедневным инструментом, а главный экран должен быстро показывать финансовое состояние и снижать риск ошибок при действиях с деньгами. "
    "Цель работы — разработать и реализовать PWA-прототип главного экрана с персонализируемыми виджетами и базовыми сценариями: просмотр баланса, перевод по номеру телефона и история операций. "
    "В первой главе я анализирую предметную область, UX-критерии и аналоги банковских интерфейсов. Во второй главе проектирую структуру, пользовательские сценарии, UI-kit и состояния интерфейса. "
    "ИТ-компонент реализован на React: есть PWA, LocalStorage, адаптивная верстка, история операций, настройки, QR-сканер и обработка ошибок. "
    "Предварительное юзабилити-тестирование показало, что ключевые сценарии выполнимы, а найденные риски частично исправлены. Дальше планируется внешнее тестирование N=5 и финальная вычитка работы."
)
doc.save(DOCX_OUT)

md = [
    "# Доклад к презентации: проектирование главного экрана мобильного банковского приложения",
    "",
    "**Студент:** Зуев Никита Игоревич, группа 8.209-2",
    "**Научный руководитель:** Коняева Софья Игоревна",
    "",
    "Рекомендуемый темп: 6–8 минут.",
    "",
]
for i, (title, speech) in enumerate(slides_data, start=1):
    md.append(f"## Слайд {i}. {title}")
    md.append(speech)
    md.append("")
MD_OUT.write_text("\n".join(md), encoding="utf-8")

print("PPTX", PPTX_OUT, PPTX_OUT.exists(), PPTX_OUT.stat().st_size)
print("DOCX", DOCX_OUT, DOCX_OUT.exists(), DOCX_OUT.stat().st_size)
print("MD", MD_OUT, MD_OUT.exists(), MD_OUT.stat().st_size)
